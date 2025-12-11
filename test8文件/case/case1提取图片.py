"""
提取文档中的图片

支持格式：
docx
ppt
pdf

"""
import fitz
from docx import Document
import os
from pptx import Presentation


def docxImage(docPath: str, outputFolder: str) -> list[str]:
    """
    提取docx文件中的图片

    :param docPath: docx文件
    :param outputFolder: 输出目录
    :return: 输出图片的文件路径列表
    """
    doc = Document(docPath)  # 打开文档
    if not os.path.exists(outputFolder):
        os.makedirs(outputFolder)  # 创建输出文件夹
    imageFiles = []
    # 遍历文档中的所有元素
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            imagePart = rel.target_part
            image_filename = os.path.join(outputFolder, os.path.basename(imagePart.partname))
            with open(image_filename, 'wb') as f:
                imageFiles.append(image_filename)
                f.write(imagePart.blob)
    return imageFiles


def pdfImage(pdfPath: str, outputFolder: str) -> list[str]:
    """
    提取pdf文件中的图片到outputFolder

    :param pdfPath: pdf文件路径
    :param outputFolder: 输出目录
    :return: 输出图片路径列表
    """
    pdf_document = fitz.open(pdfPath)
    if not os.path.exists(outputFolder):
        os.makedirs(outputFolder)
    imageFiles = []
    # 遍历每一页
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        image_list = page.get_images(full=True)

        # 遍历页面中的所有图片
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_filename = os.path.join(outputFolder, f"page_{page_num + 1}_img_{img_index + 1}.{image_ext}")

            with open(image_filename, 'wb') as f:
                imageFiles.append(image_filename)
                f.write(image_bytes)
    return imageFiles


def pptImage(pptPath: str, outputFolder: str) -> list[str]:
    """
    提取pptPath.ppt文件中的图片到outputFolder
    背景图片暂时无法提取

    :param pptPath: ppt文件路径
    :param outputFolder: 输出路径
    :return: 输出图片列表
    """
    prs = Presentation(pptPath)

    # 创建输出文件夹
    if not os.path.exists(outputFolder):
        os.makedirs(outputFolder)
    imageFiles = []
    # 遍历每一页
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 表示图片
                image = shape.image
                image_bytes = image.blob
                image_filename = os.path.join(outputFolder, f"slide_{slide.slide_id}_img_{shape.shape_id}.{image.ext}")

                with open(image_filename, 'wb') as f:
                    imageFiles.append(image_filename)
                    f.write(image_bytes)
    return imageFiles


if __name__ == '__main__':
    doc = r"C:\Users\V7.8.docx"
    # pdf = r'D:\code\pythonProject\PythonTest\test2office\file\pdf.pdf'
    # ppt = r'D:\code\pythonProject\PythonTest\test2office\file\ppt.pptx'
    # output = rf"{os.environ['USERPROFILE']}\Desktop"
    output=r"C:\Users\zhan\Desktop\图片"
    docxImage(doc, output)
    # pdfImage(pdf, output)
    # pptImage(ppt, output)
    pass
